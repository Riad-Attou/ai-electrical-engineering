#!/usr/bin/env python3
"""Build a native, editable PowerPoint for the AI in Electrical Engineering deck.

Content: Neural speed filtering for the motor + rod system (ML2).
Design: warm paper background, left rule, Georgia / Calibri / Consolas.

    python make_slides.py   ->   presentation.pptx
"""
import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
FIG  = HERE / "figures"
OUT  = HERE / "presentation.pptx"

# ── palette ────────────────────────────────────────────────────────────────
PAPER  = RGBColor(0xFB, 0xFA, 0xF6)
PAPER2 = RGBColor(0xF4, 0xF1, 0xE9)
INK    = RGBColor(0x1C, 0x2B, 0x36)
INKSOFT= RGBColor(0x4A, 0x5A, 0x66)
RULE   = RGBColor(0xD9, 0xD2, 0xC4)
NAVY   = RGBColor(0x1D, 0x35, 0x57)
TEAL   = RGBColor(0x2A, 0x9D, 0x8F)
TERRA  = RGBColor(0xE0, 0x7A, 0x5F)
GOLD   = RGBColor(0xC9, 0x9A, 0x2E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TEAL_TINT  = RGBColor(0xE7, 0xF2, 0xF0)
TERRA_TINT = RGBColor(0xFA, 0xEC, 0xE7)
GOLD_TINT  = RGBColor(0xFB, 0xF3, 0xDE)

SERIF, SANS, MONO = "Georgia", "Calibri", "Consolas"
ACCENT = {"teal": TEAL, "navy": NAVY, "terra": TERRA, "gold": GOLD}

SW, SH = Inches(13.333), Inches(7.5)
MX     = Inches(0.95)
RULE_X = Inches(0.62)
CW     = Inches(11.6)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]


# ── low-level helpers ──────────────────────────────────────────────────────

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAPER
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, RULE_X, Inches(0.55), Pt(1.4), Inches(6.4))
    _flat(rule, RULE)
    cap = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, RULE_X, Inches(0.55), Pt(2.4), Inches(0.62))
    _flat(cap, TEAL)
    return s


def _flat(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def box(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    tf.paragraphs[0].text = ""
    return tf


_TOKEN = re.compile(r"\*([^*]+)\*|~([^~]+)~|`([^`]+)`|([^*~`]+)")


def runs(p, text, font=SANS, color=INKSOFT, size=15, bold=False):
    for m in _TOKEN.finditer(text):
        strong, em, mono, plain = m.groups()
        r = p.add_run()
        if strong is not None:
            r.text, r.font.name, r.font.bold, r.font.color.rgb = strong, font, True, INK
            r.font.size = Pt(size)
        elif em is not None:
            r.text, r.font.name, r.font.bold, r.font.color.rgb = em, font, True, TEAL
            r.font.size = Pt(size)
        elif mono is not None:
            r.text, r.font.name, r.font.bold, r.font.color.rgb = mono, MONO, False, NAVY
            r.font.size = Pt(size * 0.92)
        else:
            r.text, r.font.name, r.font.bold, r.font.color.rgb = plain, font, bold, color
            r.font.size = Pt(size)


def para(tf, text, *, font=SANS, color=INKSOFT, size=15, bold=False,
         align=PP_ALIGN.LEFT, before=0, after=6, lh=1.12, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_before, p.space_after, p.line_spacing = Pt(before), Pt(after), lh
    runs(p, text, font=font, color=color, size=size, bold=bold)
    return p


def kicker(s, num, label, top=Inches(0.6)):
    tf = box(s, MX, top, CW, Inches(0.35))
    p = tf.paragraphs[0]; p.line_spacing = 1.0
    if num:
        r = p.add_run(); r.text = num + "  "
        r.font.name, r.font.size, r.font.bold, r.font.color.rgb = MONO, Pt(11), True, TERRA
    r = p.add_run(); r.text = label.upper()
    r.font.name, r.font.size, r.font.color.rgb = MONO, Pt(11), TEAL


def heading(s, text, top, size=30, w=CW):
    tf = box(s, MX, top, w, Inches(1.4))
    para(tf, text, font=SERIF, color=NAVY, size=size, bold=True, lh=1.02, first=True)


def subhead(s, text, l, t, w):
    tf = box(s, l, t, w, Inches(0.5))
    para(tf, text, font=SERIF, color=NAVY, size=19, bold=True, lh=1.05, first=True)


def bullets(s, items, l, t, w, h, size=15, gap=10):
    tf = box(s, l, t, w, h)
    for i, (txt, accent) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before, p.space_after, p.line_spacing = Pt(0 if i == 0 else gap), Pt(0), 1.1
        d = p.add_run(); d.text = "—  "
        d.font.name, d.font.size, d.font.bold = SANS, Pt(size), True
        d.font.color.rgb = ACCENT.get(accent, TERRA)
        runs(p, txt, font=SANS, color=INKSOFT, size=size)
    return tf


def card(s, l, t, w, h, *, fill=WHITE, accent=None, line=RULE):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    c.adjustments[0] = 0.06
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = line; c.line.width = Pt(1)
    c.shadow.inherit = False
    if accent:
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Pt(3), h)
        _flat(bar, ACCENT[accent])
    return c


def formula(s, lines, l, t, w, *, accent="terra", size=13.5, h=None):
    if h is None:
        h = Inches(0.38) + Inches(0.32) * len(lines)
    card(s, l, t, w, h, fill=PAPER2, accent=accent)
    tf = box(s, l + Inches(0.22), t + Inches(0.12), w - Inches(0.4), h - Inches(0.2),
             anchor=MSO_ANCHOR.MIDDLE)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before, p.space_after, p.line_spacing = Pt(0 if i == 0 else 4), Pt(0), 1.15
        runs(p, ln, font=MONO, color=NAVY, size=size)


def takeaway(s, bold, text, l, t, w, *, accent="teal", h=Inches(0.95)):
    tint = TEAL_TINT if accent == "teal" else (TERRA_TINT if accent == "terra" else GOLD_TINT)
    card(s, l, t, w, h, fill=tint, accent=accent, line=tint)
    tf = box(s, l + Inches(0.22), t + Inches(0.1), w - Inches(0.42), h - Inches(0.2),
             anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.line_spacing = 1.12
    r = p.add_run(); r.text = bold + "  "
    r.font.name, r.font.bold, r.font.size, r.font.color.rgb = SERIF, True, Pt(15), NAVY
    runs(p, text, font=SANS, color=INKSOFT, size=15)


def pills(s, items, t, *, h=Inches(1.9)):
    n = len(items)
    gap = Inches(0.25)
    w = (CW - gap * (n - 1)) / n
    for i, (tag, tac, title, body) in enumerate(items):
        l = MX + (w + gap) * i
        card(s, l, t, w, h, accent=tac)
        tf = box(s, l + Inches(0.22), t + Inches(0.18), w - Inches(0.4), h - Inches(0.3))
        para(tf, tag.upper(), font=MONO, color=ACCENT[tac], size=10, first=True, after=4)
        para(tf, title, font=SERIF, color=NAVY, size=16, bold=True, after=6, lh=1.0)
        para(tf, body, font=SANS, color=INKSOFT, size=12, lh=1.1)


def roadmap(s, items, t):
    gap = Inches(0.22)
    w = (CW - gap) / 2
    h = Inches(1.5)
    for i, (rn, title, body) in enumerate(items):
        r, c = divmod(i, 2)
        l = MX + (w + gap) * c
        tt = t + (h + gap) * r
        card(s, l, tt, w, h)
        nb = box(s, l + Inches(0.2), tt + Inches(0.16), Inches(0.7), Inches(1.2))
        para(nb, rn, font=SERIF, color=TERRA, size=30, bold=True, first=True, lh=1.0)
        tf = box(s, l + Inches(0.95), tt + Inches(0.16), w - Inches(1.15), h - Inches(0.3))
        para(tf, title, font=SERIF, color=NAVY, size=17, bold=True, first=True, after=4, lh=1.0)
        para(tf, body, font=SANS, color=INKSOFT, size=12, lh=1.12)


def stats(s, items, t):
    gap = Inches(0.3)
    w = (CW - gap * 2) / 3
    for i, (val, ck, label) in enumerate(items):
        l = MX + (w + gap) * i
        top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(2))
        _flat(top, RULE)
        tf = box(s, l, t + Inches(0.16), w, Inches(1.7))
        para(tf, val, font=SERIF, color=ACCENT.get(ck, NAVY), size=44, bold=True, first=True, lh=1.0)
        para(tf, label.upper(), font=MONO, color=INKSOFT, size=10, before=8, lh=1.1)


def figure(s, name, caption, l, t, w, *, max_h=Inches(4.9)):
    path = FIG / name
    iw, ih = Image.open(path).size
    pad = Inches(0.16)
    disp_w = w - pad * 2
    disp_h = int(disp_w * ih / iw)
    if disp_h > max_h:
        disp_h = max_h
        disp_w = int(disp_h * iw / ih)
    cap_h = Inches(0.5)
    card_h = disp_h + pad * 2 + cap_h
    card(s, l, t, w, card_h)
    img_l = l + (w - disp_w) // 2
    s.shapes.add_picture(str(path), img_l, t + pad, width=disp_w, height=disp_h)
    tf = box(s, l + pad, t + pad + disp_h + Inches(0.05), w - pad * 2, cap_h)
    para(tf, caption, font=MONO, color=INKSOFT, size=12, first=True, lh=1.1,
         align=PP_ALIGN.CENTER)
    return card_h


def tbl(s, rows, l, t, w):
    n = len(rows)
    gt = s.shapes.add_table(n, 2, l, t, w, Inches(0.45) * n).table
    gt.columns[0].width = int(w * 0.42)
    gt.columns[1].width = w - int(w * 0.42)
    gt.first_row = False; gt.horz_banding = False
    for ri, (k, v) in enumerate(rows):
        for ci, (txt, fnt, col, sz) in enumerate(
            [(k, MONO, NAVY, 12), (v, SANS, INKSOFT, 13)]
        ):
            cell = gt.cell(ri, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = PAPER
            cell.margin_left = Inches(0.08); cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            runs(p, txt, font=fnt, color=col, size=sz)


# ── slides ─────────────────────────────────────────────────────────────────

def s02_toc():
    s = slide()
    kicker(s, "", "Table of contents")
    heading(s, "Outline.", Inches(1.02), size=28)
    roadmap(s, [
        ("I",   "Physical setup & dataset",
         "Brushed DC motor + vertical rod load; excitation & sensor-noise design."),
        ("II",  "ML1 — Motor identification",
         "[Friend's section] Identification of motor dynamics."),
        ("III", "ML2 — Speed filtering",
         "Neural filter (GRU / CNN / TCN) vs EMA & Kalman; nonlinear-plant gain."),
        ("IV",  "Conclusion",
         "Key numbers, takeaways, and future directions."),
    ], Inches(2.0))


def s05_task_data():
    s = slide()
    kicker(s, "03", "ML2 — speed filter: task & data")
    heading(s, "Denoise the speed of a nonlinear plant.", Inches(1.05), size=28)
    tf = box(s, MX, Inches(2.0), CW, Inches(0.72))
    p = tf.paragraphs[0]; p.line_spacing = 1.2
    runs(p, "Recover the true motor speed from a noisy speed sensor and the known command voltage. The vertical rod adds a gravity torque *∝ sin θ*, making the mechanics *nonlinear and state-dependent* — the property a learned filter can exploit.",
         font=SANS, color=INKSOFT, size=15)
    formula(s, [
        "(J + I_rod)·dω/dt  =  Kt·i − B·ω − m·g·l_cm·sin(θ)     ← gravity = nonlinearity",
        "input   x = [ omega_noisy , voltage ]        target   y = omega_true",
    ], MX, Inches(3.05), CW, accent="navy", h=Inches(0.98))
    bullets(s, [
        ("*210 trajectories*, 6 s @ 1 ms, split *70/15/15 by trajectory* — no time-step leaks between train and test.", "teal"),
        ("Excitation cycles *step / ramp / random / mixed* bipolar voltage; motor R, J, B jittered ±12% per run.", "navy"),
        ("Sensor noise std 2–4 rad/s on a ~21 rad/s swing → *SNR ≈ 7× (14% noise)*.", "terra"),
        ("A *chirp* excitation is held out entirely for an out-of-distribution generalisation test.", "gold"),
    ], MX, Inches(4.35), CW, Inches(2.6), size=15, gap=11)


def s06_methods():
    s = slide()
    kicker(s, "04", "ML2 — methods")
    heading(s, "Classical baselines vs learned filters.", Inches(1.05), size=28)
    tf = box(s, MX, Inches(2.0), CW, Inches(0.66))
    p = tf.paragraphs[0]; p.line_spacing = 1.2
    runs(p, "Baselines tuned on validation, evaluated on test: *EMA* (α tuned) and a steady-state *Kalman* filter on the nominal linear `[i, ω]` model — whose model *omits the gravity term*. Three learned filters share the same causal *64 ms* window over [omega_noisy, voltage].",
         font=SANS, color=INKSOFT, size=15)
    pills(s, [
        ("gru", "teal",  "GRU  (3 489 params)",
         "Recurrent hidden state accumulates *64 ms* of context. hidden = 32, 1 layer."),
        ("cnn", "navy",  "CNN  (8 801 params)",
         "Stacked causal 1-D convolutions. channels = 32, k = 8, depth = 2."),
        ("tcn", "terra", "TCN  (33 153 params)",
         "Dilated causal convolutions. Receptive field = *91 ms*. channels = 32, k = 4, 4 levels."),
    ], Inches(2.95), h=Inches(2.05))
    formula(s, [
        "all learned filters:   x = (B, 64, 2)  →  scalar omega_true",
        "loss = MSE        optimiser = Adam (lr 1e-3)        early stop on val",
    ], MX, Inches(5.35), CW, accent="navy", h=Inches(0.98))


def s07_results():
    s = slide()
    kicker(s, "05", "ML2 — results  (in-distribution test set)")
    heading(s, "Learned filters cut error ~51% below Kalman.", Inches(1.05), size=27)
    tbl(s, [
        ("Raw  —  no filter",              "RMSE  2.97 rad/s   (28.3 RPM)    —   reference"),
        ("MA  —  window 64",               "RMSE  3.00 rad/s   (28.6 RPM)   −1 %   (just lag)"),
        ("Kalman  —  tuned, linear model", "RMSE  1.32 rad/s   (12.6 RPM)   +56 %"),
        ("EMA  —  tuned α",                "RMSE  1.07 rad/s   (10.2 RPM)   +64 %"),
        ("CNN",                            "RMSE  0.80 rad/s   ( 7.6 RPM)   +73 %"),
        ("*GRU*  /  *TCN 🏆*",             "RMSE  *0.65 rad/s*   ( 6.2 RPM)   *+78 %*"),
    ], MX, Inches(2.2), CW)
    takeaway(s, "Model-based ≠ better.",
             "The tuned linear *Kalman (1.32)* actually loses to a plain tuned *EMA (1.07)*: its linear model can't represent the rod's gravity torque, so its model-mismatch outweighs its model-based advantage.",
             MX, Inches(5.35), CW, accent="terra", h=Inches(1.1))


def s08_comparison():
    s = slide()
    kicker(s, "06", "ML2 — comparison")
    heading(s, "Every neural filter beats every classical one.", Inches(1.05), size=26)
    bullets(s, [
        ("A 64-sample *moving average barely helps* — the rod swings fast enough that its lag cancels its smoothing.", "navy"),
        ("*EMA and Kalman* land at 1.1–1.3 rad/s — a useful but limited linear fit.", "terra"),
        ("*GRU / TCN reach 0.65 rad/s* — a 78% cut from the raw sensor and ~51% below the tuned Kalman.", "teal"),
    ], MX, Inches(2.2), Inches(3.95), Inches(3.0), size=15, gap=14)
    figure(s, "comparison_rmse.png",
           "Fig. 1 — test-set RMSE by method (motor + rod, lower is better).",
           Inches(5.15), Inches(1.95), Inches(7.23))


def s09_generalization():
    s = slide()
    kicker(s, "07", "ML2 — generalisation  (unseen excitation)")
    heading(s, "Classical filters break on unseen inputs.", Inches(1.05), size=26)
    tbl(s, [
        ("Raw  —  no filter",      "3.00 rad/s      —     reference"),
        ("EMA  —  tuned",          "2.64 rad/s   +12 %   (was +64 %)"),
        ("Kalman  —  tuned",       "2.29 rad/s   +24 %   (was +56 %)"),
        ("TCN",                    "0.98 rad/s   +67 %   holds"),
        ("GRU",                    "0.93 rad/s   +69 %   holds"),
        ("*CNN 🏆*",               "*0.81 rad/s*   *+73 %*   holds"),
    ], MX, Inches(2.25), Inches(4.95))
    takeaway(s, "Generalisation gap.",
             "On a *chirp* sweep never seen in training, *EMA and Kalman collapse* — implicitly tuned to the training spectrum — while *neural filters hold +67–73 %*, having learned the dynamics, not the excitation.",
             MX, Inches(5.1), Inches(4.95), accent="teal", h=Inches(1.7))
    figure(s, "ood_rmse.png",
           "Fig. 2 — OOD (chirp) RMSE: classical filters collapse, neural hold.",
           Inches(6.2), Inches(1.95), Inches(6.18))


# ── main ───────────────────────────────────────────────────────────────────

def main():
    for fn in [
        s02_toc,
        s05_task_data,
        s06_methods,
        s07_results,
        s08_comparison,
        s09_generalization,
    ]:
        fn()
    prs.save(OUT)
    print(f"Wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
